from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text
    
    # Simple styling
    title.text_frame.paragraphs[0].font.bold = True
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)

def add_bullet_slide(prs, title_text, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = title_text
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = bullets[0]
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0

def main():
    prs = Presentation()
    
    # Slide 1: Title
    add_title_slide(prs, 
                   "MissionOS", 
                   "AI Mission Control for Impossible Deadlines\nTeam VibeHackathon")
    
    # Slide 2: Problem
    add_bullet_slide(prs, 
                    "The Problem: Sunk-Cost Fallacy",
                    [
                        "At the 4-hour mark of any hackathon:",
                        "- Teams realize they can't finish everything.",
                        "- But they refuse to cut their favorite features.",
                        "- They crunch, panic, and submit a broken mess.",
                        "Current tools (like Jira) only track tasks. They are passive. They don't intervene when you are about to fail."
                    ])

    # Slide 3: Solution
    add_bullet_slide(prs, 
                    "The Solution: Active Mathematical Triage",
                    [
                        "MissionOS is an active intervention agent.",
                        "- Reads Reality: Scrapes your GitHub repo (issues, PRs, code).",
                        "- Proves Failure: Uses Gemini 1.5 Flash to calculate exact failure probability.",
                        "- Forces Cuts: Outputs a brutal, prioritized list of scope cuts you MUST make to survive."
                    ])
                    
    # Slide 4: Secret Weapon
    add_bullet_slide(prs, 
                    "The Secret Weapon: Drift Detection",
                    [
                        "What happens when you accept the cuts, but then miss the new deadline?",
                        "- Timeline Monitoring: The agent watches your progress.",
                        "- Autonomous Renegotiation: If you finish a milestone late, MissionOS instantly jumps in.",
                        "- It recalculates the math and slashes even more features without you asking.",
                        "It is a true, autonomous triage agent."
                    ])

    # Slide 5: Why We Win
    add_bullet_slide(prs, 
                    "Why We Will Win",
                    [
                        "- Subverts the Prompt: Not a generic 'friendly AI to-do list', but a ruthless, high-stakes emergency tool.",
                        "- Agentic Depth: Interacts directly with live GitHub environments and actively intervenes on schedule drifts.",
                        "- Perfect Pitch: Skipped the dashboard slop. Built one continuous 'Terminal Noir' narrative flow."
                    ])

    prs.save('MissionOS_Pitch_Deck.pptx')
    print("Successfully generated MissionOS_Pitch_Deck.pptx")

if __name__ == "__main__":
    main()
